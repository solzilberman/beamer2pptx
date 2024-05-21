import subprocess
import os 
import sys 
import shutil
from pptx import Presentation
from pptx.util import Inches
import glob
import tqdm
import argparse

_TMP_DIR="_imgs"

def check_pdftoppm():
    _path = shutil.which("pdftoppm")
    if _path is None:
        raise FileNotFoundError("pdftoppm not found. Please install poppler-utils\n\tsudo apt install poppler-utils")
        
def run(infile, outfile):
    if not os.path.exists(infile):
        raise FileNotFoundError(f"{infile} not found")
    if not infile.endswith(".pdf"):
        raise ValueError("input file must be a pdf file")
    check_pdftoppm()

    print(f"[INFO] Running pdftoppm...")
    cmd=f"pdftoppm {infile} {_TMP_DIR}/slide -png -scale-to-x 1920 -scale-to-y 1080"
    if os.path.exists(_TMP_DIR):
        shutil.rmtree(_TMP_DIR)
    os.mkdir(_TMP_DIR)
    subprocess.run(cmd, shell=True)

    if os.path.exists(outfile):
        os.remove(outfile)
    slide_pngs = glob.glob(f"{_TMP_DIR}/*.png")
    slide_pngs.sort()

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    print(f"[INFO] Creating pptx file...")
    slide_iter = tqdm.tqdm(enumerate(slide_pngs), total=len(slide_pngs))
    for idx, slide_png in slide_iter:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(0)
        pic = slide.shapes.add_picture(slide_png, left, top, width=prs.slide_width, height=prs.slide_height)
    prs.save(outfile)
    print(f"[INFO] pptx file saved to {outfile}")


def parse_args():
    parser = argparse.ArgumentParser(description="convert beamer pdf to pptx")
    parser.add_argument("-i", "--infile", required=True, type=str, help="beamer pdf file name")
    parser.add_argument("-o", "--outfile", type=str, default=None, help="output pptx file name")
    return parser.parse_args()

if __name__ == "__main__":
    args = parse_args()
    if args.outfile is None:
        args.outfile = args.infile.replace(".pdf", ".pptx")
    run(args.infile, args.outfile)